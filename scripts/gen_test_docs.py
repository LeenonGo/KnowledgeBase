#!/usr/bin/env python3
"""生成测试文档：Excel、CSV、PPT"""
import os
from pathlib import Path

OUT_DIR = Path("/home/lee/.openclaw/workspace/knowledge-base/data/test_docs")
OUT_DIR.mkdir(parents=True, exist_ok=True)

# ═══════════════════════════════════════════
# 测试文档 1: 员工手册.xlsx（3 个 Sheet）
# ═══════════════════════════════════════════
def create_employee_handbook():
    from openpyxl import Workbook
    wb = Workbook()

    # Sheet 1: 考勤制度
    ws1 = wb.active
    ws1.title = "考勤制度"
    ws1.append(["条款编号", "条款名称", "详细内容", "生效日期"])
    ws1.append(["ATT-001", "工作时间", "公司实行标准工时制，周一至周五 9:00-18:00，午休 12:00-13:00。研发部门可弹性到 10:00-19:00。", "2024-01-01"])
    ws1.append(["ATT-002", "打卡规定", "每日上下班需在钉钉打卡，迟到 30 分钟以内扣半天绩效，超过 30 分钟按旷工半天处理。", "2024-01-01"])
    ws1.append(["ATT-003", "加班管理", "加班需提前在 OA 系统提交申请，经直属领导审批后方可生效。工作日加班按 1.5 倍调休，周末 2 倍，法定节假日 3 倍。", "2024-03-01"])
    ws1.append(["ATT-004", "年假制度", "入职满 1 年享 5 天年假，满 3 年享 10 天，满 5 年享 15 天。年假须在当年度内休完，不可跨年累计。", "2024-01-01"])
    ws1.append(["ATT-005", "病假规定", "病假需提供三甲医院诊断证明，3 天以内直属领导审批，超过 3 天需 HR 总监审批。病假期间发放基本工资的 80%。", "2024-06-01"])
    ws1.append(["ATT-006", "远程办公", "经部门负责人批准，每周可申请 1-2 天远程办公。远程期间需保持钉钉在线，参加线上会议。", "2024-09-01"])

    # Sheet 2: 报销流程
    ws2 = wb.create_sheet("报销流程")
    ws2.append(["流程步骤", "负责人", "操作说明", "时限要求"])
    ws2.append(["Step 1: 提交申请", "报销人", "在 OA 系统填写报销单，上传发票照片（增值税专用发票需提供完整联次），选择费用类别（差旅/招待/办公/培训）。", "费用发生后 30 天内"])
    ws2.append(["Step 2: 部门审核", "直属领导", "核实费用真实性与合理性，确认是否在部门预算范围内。差旅报销需附出差审批单。", "收到后 2 个工作日"])
    ws2.append(["Step 3: 财务复核", "财务专员", "审核发票合规性、金额准确性、预算科目匹配度。不符合规定的退回并注明原因。", "收到后 3 个工作日"])
    ws2.append(["Step 4: 审批", "财务总监", "单笔 5000 元以下自动通过，5000-20000 元需财务总监审批，20000 以上需总经理审批。", "收到后 2 个工作日"])
    ws2.append(["Step 5: 付款", "出纳", "审批通过后 5 个工作日内打款至员工工资卡。如有借款需先冲抵。", "审批后 5 个工作日"])
    ws2.append(["特殊情况", "—", "发票丢失需提交情况说明并由部门负责人签字确认，按发票金额 80% 报销。跨年发票原则上不予报销。", "—"])

    # Sheet 3: 绩效考核
    ws3 = wb.create_sheet("绩效考核")
    ws3.append(["考核维度", "权重", "评分标准", "数据来源"])
    ws3.append(["工作成果", "40%", "按时交付率 ≥ 95% 满分，90-95% 扣 10%，<90% 扣 30%。项目延期需提前报备，未报备直接扣分。", "项目管理系统"])
    ws3.append(["代码质量", "20%", "Bug 率 < 5% 满分，5-10% 扣 20%，>10% 扣 50%。Code Review 通过率 ≥ 90%。", "GitLab / 禅道"])
    ws3.append(["团队协作", "15%", "360 度评价均分。主动协助他人解决问题额外加分。跨部门协作满意度调查。", "360 评价系统"])
    ws3.append(["技术创新", "15%", "技术分享每次 +5 分，专利申报 +10 分，开源贡献 +5 分，内部工具优化 +3 分。", "技术委员会记录"])
    ws3.append(["学习成长", "10%", "完成年度培训计划满分，获得行业认证额外 +5 分，参加外部技术会议 +3 分。", "培训系统"])

    wb.save(str(OUT_DIR / "员工手册_考勤报销绩效.xlsx"))
    print("✅ 员工手册_考勤报销绩效.xlsx")


# ═══════════════════════════════════════════
# 测试文档 2: 项目台账.csv
# ═══════════════════════════════════════════
def create_project_ledger():
    import csv
    rows = [
        ["项目编号", "项目名称", "负责人", "起始日期", "截止日期", "当前状态", "预算(万元)", "完成度", "备注"],
        ["PRJ-2024-001", "智能客服系统升级", "张伟", "2024-01-15", "2024-06-30", "已完成", "85.0", "100%", "上线后用户满意度提升 23%"],
        ["PRJ-2024-002", "数据中台建设", "李娜", "2024-03-01", "2024-12-31", "进行中", "200.0", "65%", "数据治理模块已完成，可视化看板开发中"],
        ["PRJ-2024-003", "移动端 App 重构", "王强", "2024-04-15", "2024-09-30", "进行中", "120.0", "40%", "Flutter 技术栈迁移，iOS 端已提交审核"],
        ["PRJ-2024-004", "安全合规审计", "赵敏", "2024-06-01", "2024-08-31", "已暂停", "30.0", "25%", "等待等保 2.0 测评机构排期"],
        ["PRJ-2024-005", "RAG 知识库系统", "陈宇", "2024-07-01", "2025-01-31", "进行中", "95.0", "80%", "已完成 OCR 模块，正在做效果评测"],
        ["PRJ-2024-006", "CI/CD 流水线优化", "刘洋", "2024-08-01", "2024-10-31", "已完成", "15.0", "100%", "构建时间从 15 分钟降到 4 分钟"],
        ["PRJ-2024-007", "用户画像平台", "孙丽", "2024-09-01", "2025-03-31", "规划中", "150.0", "5%", "需求评审阶段，已完成用户调研"],
        ["PRJ-2024-008", "API 网关升级", "周杰", "2024-10-15", "2025-01-15", "进行中", "45.0", "55%", "限流熔断已上线，灰度发布功能开发中"],
    ]

    with open(str(OUT_DIR / "项目台账_2024.csv"), "w", newline="", encoding="utf-8-sig") as f:
        writer = csv.writer(f)
        writer.writerows(rows)
    print("✅ 项目台账_2024.csv")


# ═══════════════════════════════════════════
# 测试文档 3: 产品汇报.pptx（5 页）
# ═══════════════════════════════════════════
def create_product_report():
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide(title, body_lines, notes=""):
        layout = prs.slide_layouts[1]  # 标题+内容
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        for i, line in enumerate(body_lines):
            if i == 0:
                tf.text = line
            else:
                p = tf.add_paragraph()
                p.text = line
        if notes and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = notes

    # 第 1 页：封面
    add_slide(
        "RAG 知识库管理系统 — 2024 年度汇报",
        ["汇报人：产品团队", "日期：2025 年 1 月 15 日", "版本：v3.0"],
        notes="准备时间 15 分钟，重点讲 Phase 4-7 的进展"
    )

    # 第 2 页：项目里程碑
    add_slide(
        "项目里程碑",
        [
            "Phase 1-2（Q1）：MVP + 文档管理 — 3 月上线内部测试",
            "Phase 3（Q2）：安全架构 + 混合检索 — 通过等保预审",
            "Phase 4-5（Q3）：质量监控 + 查询改写 + Reranker — 差评率从 18% 降到 6%",
            "Phase 6-7（Q4）：效果评测 + OCR — 支持 PDF 自动解析，评测覆盖 9 维度",
        ],
        notes="强调迭代速度：4 个 Phase，7 个月从 0 到生产可用"
    )

    # 第 3 页：核心指标
    add_slide(
        "核心业务指标",
        [
            "日均问答量：从上线初 50 次增长到 800+ 次",
            "平均响应时间：< 2 秒（含 Reranker）",
            "用户满意度：92%（好评 / 总反馈）",
            "检索精确率：87%（Reranker 优化后提升 15%）",
            "支持文档格式：PDF / Word / TXT / Markdown，共入库 1,200+ 份文档",
            "知识库数量：12 个部门知识库，3 个公共库",
        ],
        notes="数据截至 12 月 31 日，来自系统内置质量监控模块"
    )

    # 第 4 页：技术架构
    add_slide(
        "技术架构亮点",
        [
            "混合检索：向量语义 + BM25 关键词 + RRF 融合，召回率提升 30%",
            "重排序：qwen3-vl-rerank 二次排序，过滤低质量检索结果",
            "查询改写：多轮对话指代消解，上下文理解准确率 95%",
            "效果评测：LLM-as-Judge 自动评分，9 维度持续监控",
            "OCR 引擎：PaddleOCR 版面检测 + 文字识别 + 表格识别",
            "权限隔离：部门级授权 + 三级角色权限，检索时自动过滤",
        ],
        notes="技术选型理由：ChromaDB 轻量适合单机部署，DashScope 兼容 OpenAI 接口方便切换"
    )

    # 第 5 页：下一步计划
    add_slide(
        "2025 Q1-Q2 规划",
        [
            "1. Query 润色：拼写纠错 + 同义扩展 + 关键词提取（预计 1 月完成）",
            "2. 多格式支持：Excel / CSV / PPT 文档解析（预计 2 月完成）",
            "3. 数据源同步：飞书 / Confluence 自动导入（预计 3 月完成）",
            "4. Redis 缓存：替换内存缓存，支持多实例部署（预计 4 月完成）",
            "5. API 开放 + Bot 发布：对外提供 API 和嵌入式 Widget（预计 5 月完成）",
            "",
            "目标：用户量从 200 人扩展到 1000 人，文档量突破 10,000 份",
        ],
        notes="每个功能点需要 2-4 周开发周期，按优先级排序"
    )

    prs.save(str(OUT_DIR / "产品年度汇报_2024.pptx"))
    print("✅ 产品年度汇报_2024.pptx")


# ═══════════════════════════════════════════
# 测试文档 4: 财务数据.xlsx（纯数据表）
# ═══════════════════════════════════════════
def create_finance_data():
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "季度收入"
    ws.append(["季度", "产品线", "收入(万元)", "成本(万元)", "利润率", "同比增长"])
    data = [
        ["Q1", "SaaS 订阅", 320, 85, "73.4%", "+25%"],
        ["Q1", "定制开发", 180, 120, "33.3%", "-10%"],
        ["Q1", "技术支持", 95, 30, "68.4%", "+15%"],
        ["Q2", "SaaS 订阅", 380, 90, "76.3%", "+32%"],
        ["Q2", "定制开发", 150, 95, "36.7%", "-5%"],
        ["Q2", "技术支持", 110, 35, "68.2%", "+20%"],
        ["Q3", "SaaS 订阅", 450, 95, "78.9%", "+40%"],
        ["Q3", "定制开发", 200, 130, "35.0%", "+5%"],
        ["Q3", "技术支持", 125, 40, "68.0%", "+28%"],
        ["Q4", "SaaS 订阅", 520, 100, "80.8%", "+48%"],
        ["Q4", "定制开发", 160, 100, "37.5%", "0%"],
        ["Q4", "技术支持", 140, 45, "67.9%", "+35%"],
    ]
    for row in data:
        ws.append(row)
    wb.save(str(OUT_DIR / "财务数据_2024季度收入.xlsx"))
    print("✅ 财务数据_2024季度收入.xlsx")


# ═══════════════════════════════════════════
# 测试文档 5: 培训资料.pptx（含备注）
# ═══════════════════════════════════════════
def create_training_material():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    def add_slide(title, bullets, notes=""):
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        for i, b in enumerate(bullets):
            if i == 0:
                tf.text = b
            else:
                p = tf.add_paragraph()
                p.text = b
        if notes and slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.text = notes

    add_slide(
        "新员工入职培训 — 技术篇",
        ["欢迎加入研发团队！", "本次培训涵盖：开发环境搭建、代码规范、CI/CD 流程"],
        notes="培训时长约 2 小时，预留 30 分钟 Q&A"
    )
    add_slide(
        "开发环境搭建",
        [
            "1. 安装 Docker Desktop（Mac/Windows）或 Docker Engine（Linux）",
            "2. 克隆项目仓库：git clone git@gitlab.company.com:team/project.git",
            "3. 安装依赖：docker-compose up -d（启动 MySQL、Redis、ES）",
            "4. 配置 IDE：推荐 VS Code + Remote Container 插件",
            "5. 运行测试：make test（确保本地环境正常）",
        ],
        notes="常见问题：Docker 镜像拉取慢可配置公司镜像源，地址在 Wiki 上"
    )
    add_slide(
        "代码规范",
        [
            "Python: 遵循 PEP 8，使用 black 格式化，pylint 评分 ≥ 8.0",
            "JavaScript: ESLint + Prettier，React 项目用 Airbnb 规范",
            "Git 提交: 遵循 Conventional Commits（feat/fix/docs/refactor）",
            "Code Review: 每个 MR 至少 2 人 Review，CI 通过后方可合并",
            "文档: 公共 API 必须写 docstring，复杂逻辑加行内注释",
        ],
        notes="代码规范检查已集成到 CI 流水线，不合规的 MR 无法合并"
    )
    add_slide(
        "CI/CD 流程",
        [
            "提交代码 → GitLab CI 自动触发",
            "Stage 1: 代码检查（lint + type check + security scan）",
            "Stage 2: 单元测试（pytest + coverage ≥ 80%）",
            "Stage 3: 构建镜像（Docker multi-stage build）",
            "Stage 4: 部署到 Staging（自动）→ 手动确认 → 部署到 Production",
            "回滚：kubectl rollout undo 或 CI 面板一键回滚",
        ],
        notes="生产部署窗口：每周二、四的 14:00-16:00，紧急 hotfix 除外"
    )
    add_slide(
        "常见问题 & 联系方式",
        [
            "Q: Git 权限被拒绝？→ 检查 SSH key 是否添加到 GitLab（Settings → SSH Keys）",
            "Q: Docker 容器启动失败？→ 查看 docker-compose logs，常见原因是端口冲突",
            "Q: 测试数据库连不上？→ 确认 .env 中 DB_HOST 指向 docker 网络地址",
            "",
            "技术支持：tech-support@company.com",
            "IT 热线：ext. 8888",
        ],
        notes="培训录屏会上传到内部学习平台，方便回看"
    )

    prs.save(str(OUT_DIR / "新员工培训_技术篇.pptx"))
    print("✅ 新员工培训_技术篇.pptx")


if __name__ == "__main__":
    create_employee_handbook()
    create_project_ledger()
    create_product_report()
    create_finance_data()
    create_training_material()
    print(f"\n📁 所有测试文档已生成到: {OUT_DIR}")
