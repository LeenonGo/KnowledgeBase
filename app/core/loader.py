"""文档加载器 — 支持 TXT、Word、PDF、Markdown"""

import logging
from pathlib import Path

logger = logging.getLogger("kb.loader")


def load_document(file_path: str, progress_callback=None) -> tuple[str, list[str]]:
    """
    根据文件扩展名选择加载方式，返回 (纯文本, 警告列表)。
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    loaders = {
        ".txt": _load_txt,
        ".md": _load_txt,
        ".pdf": _load_pdf,
        ".docx": _load_docx,
        ".xlsx": _load_excel,
        ".xls": _load_excel,
        ".csv": _load_csv,
        ".pptx": _load_pptx,
    }

    loader = loaders.get(ext)
    if loader is None:
        raise ValueError(f"不支持的文件格式: {ext}")

    if ext == ".pdf":
        return loader(path, progress_callback=progress_callback)
    return loader(path)


def _load_txt(path: Path) -> tuple[str, list[str]]:
    """加载 TXT / Markdown"""
    return path.read_text(encoding="utf-8"), []


def _load_pdf(path: Path, progress_callback=None) -> tuple[str, list[str]]:
    """加载 PDF：统一走 OCR 解析"""
    warnings = []
    logger.info("OCR 解析 PDF: %s", path.name)

    try:
        from app.core.ocr import OCREngine, build_clean_output, build_markdown

        engine = OCREngine(device="cpu")
        raw = engine.analyze_pdf(str(path), progress_callback=progress_callback)

        # 收集引擎处理中的错误
        for page in raw.get("pages", []):
            for region in page.get("regions", []):
                if region.get("error"):
                    label = region.get("type", "unknown")
                    warnings.append(f"区域 [{label}] 处理失败: {region['error']}")

        clean = build_clean_output(raw)
        md = build_markdown(clean)

        if md.strip():
            logger.info("OCR 完成，提取 %d 字符", len(md.strip()))
            return md.strip(), warnings
    except ValueError:
        # 预检失败（模型/依赖缺失）— 直接向上抛，不 fallback
        raise
    except Exception as e:
        msg = f"OCR 处理失败: {e}"
        logger.warning(msg)
        warnings.append(msg)

    # OCR 处理异常（非预检问题），回退 PyMuPDF
    import fitz
    doc = fitz.open(str(path))
    text = "\n".join(page.get_text() for page in doc)
    doc.close()
    warnings.append("已回退到基础文本提取（PyMuPDF），内容可能不完整")
    return text.strip(), warnings


def _load_docx(path: Path) -> tuple[str, list[str]]:
    """加载 Word 文档"""
    from docx import Document

    doc = Document(str(path))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return text, []


def _excel_rows_to_records(ws, sheet_name: str = "") -> list[str]:
    """将 Excel worksheet 转为每条记录独立的 Markdown 段落（每条记录一个 ## 标题）"""
    rows_data = []
    header = []
    for i, row in enumerate(ws.iter_rows(values_only=True)) if hasattr(ws, 'iter_rows') else enumerate(range(ws.nrows)):
        if hasattr(ws, 'iter_rows'):
            # openpyxl
            if i == 0:
                header = [str(c) if c is not None else f"列{k+1}" for k, c in enumerate(row)]
                continue
            if all(c is None for c in row):
                continue
            cells = list(row)
        else:
            # xlrd
            cells = [ws.cell_value(i, j) for j in range(ws.ncols)]
            if i == 0:
                header = [str(c) if c else f"列{j+1}" for j, c in enumerate(cells)]
                continue
            if all(str(c).strip() == "" for c in cells):
                continue

        # 取第一个非编号字段做标题
        title = ""
        parts = []
        for j, val in enumerate(cells):
            is_none = val is None if hasattr(ws, 'iter_rows') else str(val).strip() == ""
            if is_none:
                continue
            col_name = header[j] if j < len(header) else f"列{j+1}"
            val_str = str(val)
            parts.append(f"{col_name}为{val_str}。")
            if not title and "编号" not in col_name and "序号" not in col_name:
                title = f"{val_str}"[:20]

        if parts:
            prefix = f"{sheet_name} - " if sheet_name else ""
            rec_title = f"{prefix}{title}" if title else f"{prefix}第 {i} 条"
            rows_data.append(f"## {rec_title}\n\n" + "\n".join(parts))
    return rows_data


def _load_excel(path: Path) -> tuple[str, list[str]]:
    """加载 Excel (.xlsx / .xls)，每条记录用 ## 标题分隔"""
    warnings = []
    ext = path.suffix.lower()

    sheets = {}
    if ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            records = _excel_rows_to_records(ws, sheet_name=sheet_name)
            if records:
                sheets[sheet_name] = records
        wb.close()
    else:  # .xls
        import xlrd
        wb = xlrd.open_workbook(str(path))
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            records = _excel_rows_to_records(ws, sheet_name=sheet_name)
            if records:
                sheets[sheet_name] = records

    if not sheets:
        warnings.append("Excel 文件为空或无法解析")
        return "", warnings

    sections = []
    for sheet_name, rows in sheets.items():
        sections.append("\n\n".join(rows))

    text = "\n\n".join(sections)
    return text, warnings


def _load_csv(path: Path) -> tuple[str, list[str]]:
    """加载 CSV，自动检测编码"""
    import csv
    warnings = []

    # 编码检测
    raw = path.read_bytes()
    encoding = "utf-8"

    try:
        import chardet
        detected = chardet.detect(raw[:10000])
        if detected.get("encoding") and detected["confidence"] > 0.5:
            encoding = detected["encoding"]
    except ImportError:
        pass

    # 尝试读取
    text = None
    for enc in [encoding, "utf-8", "gbk", "gb2312", "latin-1"]:
        try:
            text = raw.decode(enc)
            break
        except (UnicodeDecodeError, LookupError):
            continue

    if text is None:
        warnings.append("CSV 编码检测失败，使用 latin-1 兜底")
        text = raw.decode("latin-1")

    # 解析 CSV
    lines = text.strip().split("\n")
    if not lines:
        return "", ["CSV 文件为空"]

    reader = csv.reader(lines)
    rows_list = list(reader)

    if len(rows_list) < 2:
        warnings.append("CSV 仅含表头或为空")
        return rows_list[0][0] if rows_list else "", warnings

    header = rows_list[0]
    data_rows = rows_list[1:]

    output_lines = []
    for i, row in enumerate(data_rows, 1):
        parts = []
        for j, val in enumerate(row):
            if val.strip():
                col_name = header[j] if j < len(header) else f"列{j+1}"
                parts.append(f"{col_name}为{val.strip()}。")
        if parts:
            output_lines.append(f"第 {i} 条记录：" + "\n".join(parts))

    if not output_lines:
        warnings.append("CSV 无有效数据行")
        return "", warnings

    return "\n\n".join(output_lines), warnings


def _load_pptx(path: Path) -> tuple[str, list[str]]:
    """加载 PowerPoint (.pptx)，提取标题+正文+备注"""
    from pptx import Presentation
    warnings = []

    prs = Presentation(str(path))
    slides_text = []

    for i, slide in enumerate(prs.slides, 1):
        parts = []

        # 提取标题
        if slide.shapes.title and slide.shapes.title.text.strip():
            parts.append(f"标题：{slide.shapes.title.text.strip()}")

        # 提取正文文本框
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)

        # 提取备注
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"备注：{notes}")

        if parts:
            slides_text.append(f"## 幻灯片 {i}\n\n" + "\n".join(parts))
        else:
            warnings.append(f"幻灯片 {i} 无文本内容（可能含图片）")

    if not slides_text:
        warnings.append("PPT 文件为空或无文本内容")
        return "", warnings

    return "\n\n".join(slides_text), warnings
